"""
PowerPoint Generator Tool: Creates PowerPoint (.pptx) files using python-pptx.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from backend.tools.path_safety import safe_resolve_output_path

logger = logging.getLogger(__name__)

# Output directory for generated presentations
OUTPUT_DIR = Path(__file__).parent.parent.parent / "workspace" / "outputs"


def generate_ppt(filename: str, title: str, bullet_points: list) -> str:
    """
    Generate a PowerPoint presentation with a title slide and bullet points.

    Args:
        filename: Name of the output file (e.g. "slides.pptx").
        title: The slide title text.
        bullet_points: List of strings to render as bullet points.

    Returns:
        Absolute path to the created file.
    """
    # Ensure output directory exists
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Resolve path with containment check
    output_path = safe_resolve_output_path(filename, OUTPUT_DIR)

    try:
        prs = Presentation()

        # Use the blank slide layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title textbox (upper portion of the slide)
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9.0)
        height = Inches(1.0)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_tf.text = title

        # Format title text
        title_para = title_tf.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Add content textbox (below title)
        top = Inches(1.5)
        height = Inches(5.0)

        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_tf = content_box.text_frame
        content_tf.word_wrap = True

        # Add bullet points
        for i, point in enumerate(bullet_points):
            if i == 0:
                para = content_tf.paragraphs[0]
            else:
                para = content_tf.add_paragraph()

            para.text = point
            para.font.size = Pt(18)
            para.level = 0
            # Bullet character
            para.text = f"• {point}"

        # Save the presentation
        prs.save(str(output_path))

        logger.info(f"PowerPoint generated: {output_path}")
        return str(output_path)

    except Exception as e:
        logger.error(f"PowerPoint generation failed: {e}")
        raise
